
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import os
    import re
except ImportError:
    print("Error: 'python-pptx' library is not installed.")
    print("Please install it by running: pip install python-pptx")
    exit(1)

def create_presentation():
    prs = Presentation()
    
    # Read the markdown content
    with open('PRESENTATION_CONTENT.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides_content = content.split('---')
    
    for slide_text in slides_content:
        if not slide_text.strip():
            continue
            
        lines = slide_text.strip().split('\n')
        
        # Determine slide layout
        # 0 = Title Slide, 1 = Title and Content
        if "Slide 1:" in lines[0]:
            slide_layout = prs.slide_layouts[0] # Title Slide
        else:
            slide_layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(slide_layout)
        
        # Parse content
        title_text = ""
        body_text_lines = []
        
        for line in lines:
            if line.startswith("**Title:**"):
                title_text = line.replace("**Title:**", "").strip()
            elif line.startswith("**Subtitle:**"):
                body_text_lines.append(line.replace("**Subtitle:**", "").strip())
            elif line.startswith("**Presented By:**"):
                body_text_lines.append("\n" + line.replace("**Presented By:**", "Presented By: ").strip())
            elif line.startswith("*   ") or line.startswith("1.  "):
                body_text_lines.append(line)
            elif line.startswith("**Before:**") or line.startswith("**After (Eco-Command):**"):
                body_text_lines.append("\n" + line)
            elif line.startswith("*(Visual Suggestion:"):
                pass # Skip visual suggestions in text, purely for file referencing if we were advanced
            
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            
        # Set Body/Subtitle
        if slide_layout == prs.slide_layouts[0]: # Title Slide
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                body_shape.text = "\n".join(body_text_lines)
        else: # Content Slide
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = "" # Clear default
                
                for line in body_text_lines:
                    p = tf.add_paragraph()
                    clean_line = line
                    
                    # Handle bullet points
                    if line.strip().startswith("*   "):
                        clean_line = line.replace("*   ", "").strip()
                        p.level = 0
                    elif line.strip().startswith("1.  "):
                        clean_line = line.replace("1.  ", "").strip()
                        p.level = 0
                    else:
                        p.level = 0
                        
                    # Bold handling (basic)
                    if "**" in clean_line:
                        parts = clean_line.split("**")
                        p.text = "" # Reset to build runs
                        for i, part in enumerate(parts):
                            run = p.add_run()
                            run.text = part
                            if i % 2 == 1: # The text inside ** **
                                run.font.bold = True
                    else:
                         p.text = clean_line

    output_file = 'Eco_Command_Presentation.pptx'
    prs.save(output_file)
    print(f"Successfully generated '{output_file}'!")

if __name__ == "__main__":
    create_presentation()
